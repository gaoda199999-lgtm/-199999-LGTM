# -*- coding: utf-8 -*-
"""
矿山探采对比分析报告生成工具 v12.0（性能优化 + Bug 修复 + 代码质量提升）
─────────────────────────────────────────────────────────────
【V12.0 优化内容】
  1. [Bug修复] save_indicators_to_db 不再跳过核心指标（p_q/z_q/p_c/z_c/p_p/z_p），
     改为排除非核心计算字段（q_diff/q_diff_pct），所有数值型指标正常入库。
  2. [Bug修复] detect_anomaly 合并为单次数据库连接、单条 SQL（ORDER BY + LIMIT 2）。
  3. [Bug修复] validate_ore_logic 循环外预创建 ColumnMatcher，循环内复用。
  4. [性能] 新增 classify_excel_file_cached，extract_duibai_data/extract_pinghua_data
     内部只打开一次 workbook，避免重复打开 Excel 3 次。
  5. [性能] render_and_save 改为先存 PNG 再用文件路径插入 Word，同一张图只渲染一次。
  6. [代码质量] 合并 6 个重复图表函数为 3 个统一函数（make_ore_chart/make_grade_chart/
     make_metal_chart），旧函数名保留为 wrapper 向后兼容。
  7. [代码质量] 提取公共表头解析逻辑 _find_header_and_parse(ws)。
  8. [代码质量] 添加关键函数的 typing 类型提示。
  9. [清理] _fig_to_buf 保留（可能外部使用），draw_comparison_* 仍返回 (fig, buf) 保持兼容，
     但 render_and_save 内部忽略 buf，直接用 fig 存 PNG。
  10. 版本号更新为 v12.0。

  所有原有功能完全保留，向后兼容，无业务逻辑删减。
"""

import os
import sys
import re
import json
import sqlite3
import argparse
import warnings
from datetime import datetime
from io import BytesIO
from collections import defaultdict
from typing import Dict, List, Optional, Tuple, Any

import matplotlib.pyplot as plt
import matplotlib.font_manager as mpl_fm
from matplotlib.ticker import AutoLocator
from matplotlib.figure import Figure
import numpy as np
from openpyxl import load_workbook
from openpyxl.worksheet.worksheet import Worksheet
from docx import Document
from docx.shared import Pt, Cm, RGBColor, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.enum.style import WD_STYLE_TYPE
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
from docx.oxml.shared import OxmlElement as OE
from pptx import Presentation
from pptx.util import Inches as PptInches, Pt as PptPt, Emu
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.text import PP_ALIGN

# ═══════════════════════════════════════════════════════════
# 0. 全局同义词库（核心智能识别引擎）
# ═══════════════════════════════════════════════════════════

SYNONYM_MAP = {
    "生产_矿石量": ["生产矿石量", "生产矿量", "采出矿量", "采出量", "生产量",
                     "矿石量(生产)", "采出矿量(t)", "prod_ore", "ProdOre"],
    "ZHTK_矿石量": ["ZHTK矿石量", "ZHTK矿量", "地质矿量", "地质矿石量",
                     "模型矿量", "模型矿石量", "zhtk_ore", "ZhtkOre"],
    "生产_品位": ["生产品位", "生产品位(%)", "品位(生产)", "Cu(生产)",
                  "生产Cu(%)", "prod_grade", "ProdCu"],
    "ZHTK_品位": ["ZHTK品位", "ZHTK品位(%)", "地质品位", "模型品位",
                  "Cu(ZHTK)", "ZHTK Cu(%)", "zhtk_grade", "ZhtkCu"],
    "生产_金属量": ["生产金属量", "金属量(生产)", "生产铜金属量", "铜金属量(生产)",
                    "prod_metal", "ProdMetal"],
    "ZHTK_金属量": ["ZHTK金属量", "地质金属量", "模型金属量", "ZHTK铜金属量",
                    "zhtk_metal", "ZhtkMetal"],
    "面积_生产": ["生产面积", "采场面积", "面积(生产)", "生产模型面积",
                  "prod_area", "矿体面积(生产)"],
    "面积_ZHTK": ["ZHTK面积", "地质面积", "模型面积", "ZHTK模型面积",
                  "zhtk_area", "矿体面积(ZHTK)"],
    "面积_重合": ["重合面积", "重叠面积", "重合区面积", "overlap_area",
                  "面积重合"],
    "面积_歪曲增加": ["歪曲增加", "增加面积", "多圈面积", "正歪曲面积",
                      "增加区面积", "area_increase"],
    "面积_歪曲减少": ["歪曲减少", "减少面积", "少圈面积", "负歪曲面积",
                      "减少区面积", "area_decrease"],
    "贫化率": ["贫化率", "贫化(%)", "矿石贫化率", "γ", "dilution",
               "Dilution(%)"],
    "损失率": ["损失率", "损失(%)", "矿石损失率", "ρ", "loss",
               "Loss(%)"],
    "出矿量": ["出矿量", "出矿吨", "采矿量", "实际采矿量", "出矿总量",
               "采出矿量", "OreOutput"],
    "品位": ["品位", "出矿品位", "实际品位", "Cu(%)", "铜品位",
             "OreGrade"],
    "废石混入率": ["废石混入率", "混入率", "废石混入", "waste_mix"],
    "矿石总损失率": ["矿石总损失率", "总损失率", "综合损失率", "total_loss"],
    "采矿损失量": ["采矿损失量", "损失矿量", "损失矿石量", "ore_loss"],
    "贫化损失总量": ["贫化损失总量", "损失贫化合计", "贫损合计"],
}

# V11.1：表头定位关键词（合并 V10 的多关键词支持）
HEADER_ROW_KEYWORDS = ["台阶", "分段", "平台", "标高", "分台阶", "step", "bench", "level"]

def is_header_row(first_cell):
    """检查某行的第一列是否匹配表头关键词"""
    if first_cell is None:
        return False
    return any(kw in str(first_cell) for kw in HEADER_ROW_KEYWORDS)

FILE_SCORE_THRESHOLD = 0.5

# ═══════════════════════════════════════════════════════════
# 默认配置（融合 V9.1 完整配置 + V10 精简 columns）
# ═══════════════════════════════════════════════════════════

DEFAULT_CONFIG = {
    "mine": {
        "name": "米拉多铜矿",
        "subtitle": "季度探采对比分析",
        "location": "厄瓜多尔"
    },
    "files": {
        "pinghua_keywords": ["贫化损失", "统计"],
        "duibai_keywords": ["探采对比"],
        "pinghua_pattern": "{year}年矿石贫化损失统计表-上报版{month_range}.xlsx",
        "duibai_pattern": "{year}年{cn_q}季度探采对比表格.xlsx"
    },
    "columns": {
        "pinghua": {
            "grouping": "every_n",
            "n_per_step": 3,
            "skip_keywords": ["备注", "签名", "说明", "单位", "制表", "审核", "采矿场", "时间"],
            "patterns": {
                "贫化率": [["贫化", "γ"]],
                "损失率": [["损失", "ρ"]],
                "出矿量": [["出矿", "采矿量"]],
                "品位": [["品位", "Cu"]]
            },
            "required": ["贫化率", "损失率"],
            "summary_keywords": ["合计", "单位", "制表", "审核", "采矿场"]
        },
        "duibai": {
            "sheet_patterns": ["矿量对比", "面积误差"],
            "required_sheets": ["矿量对比", "面积误差重合率"],
            "patterns": {
                "生产_矿石量": [["生产"], ["矿石量", "矿量"]],
                "生产_品位": [["生产"], ["品位", "Cu"]],
                "生产_金属量": [["生产"], ["金属量"]],
                "ZHTK_矿石量": [["ZHTK"], ["矿石量", "矿量"]],
                "ZHTK_品位": [["ZHTK"], ["品位", "Cu"]],
                "ZHTK_金属量": [["ZHTK"], ["金属量"]],
                "面积_生产": [["生产", "采场"], ["面积"]],
                "面积_ZHTK": [["ZHTK", "地质"], ["面积"]],
                "面积_重合": [["重合"], ["面积"]],
                "面积_歪曲增加": [["歪曲", "增加"], ["面积"]],
                "面积_歪曲减少": [["歪曲", "减少"], ["面积"]]
            },
            "required_by_sheet": {
                "矿量对比": ["生产_矿石量", "ZHTK_矿石量", "生产_品位", "ZHTK_品位", "生产_金属量", "ZHTK_金属量"],
                "面积误差": ["面积_生产", "面积_ZHTK", "面积_重合", "面积_歪曲增加", "面积_歪曲减少"]
            },
            "summary_keywords": ["合计", "中部", "下部"]
        }
    },
    "grading": {
        "A级": {"重合率_min": 80, "歪曲率_max": 30, "误差率_max": 10},
        "B级": {"重合率_min": 80, "歪曲率_max": 40, "误差率_max": 20}
    },
    "analyze": {
        "distort_warning_pct": 25,
        "ore_error_warning_pct": 10
    },
    "report": {
        "chapter1_text": "本报告以ZHTK三维地质勘探模型（以下简称ZHTK模型）为对比基准，以同区域生产勘探模型（以下简称生产模型）为验证对象，进行季度探采对比分析。",
        "step_order": []
    },
    "styles": {
        "font_name": "宋体",
        "title_size": 20,
        "h1_size": 15,
        "h2_size": 13,
        "body_size": 11,
        "table_font_size": 8.5,
        "header_bg_color": "2F5496",
        "header_text_color": "FFFFFF",
        "body_indent_cm": 0.74
    },
    "colors": {
        "model_prod": "#2E6B8A",
        "model_zhtk": "#C98247",
        "single_blue": "#2E6B8A",
        "single_red": "#D4785A",
        "area_prod": "#3B76B8",
        "area_zhtk": "#FF901C",
        "area_overlap": "#FFC800",
        "area_rel_e": "#235E1F",
        "area_dr": "#649DC9",
        "area_wadd": "#649DC9",
        "area_wr": "#2F4F2E",
        # ===== V9 新增：歪曲减少柱独立配色 =====
        "area_decrease": "#A0A0A0"
    },
    # ===== V9 新增：柱顶标签格式统一 =====
    "bar_label_fmt": {
        "ore": "{:.1f}",
        "metal": "{:.2f}",
        "area": "{:.0f}"
    }
}

# ═══════════════════════════════════════════════════════════
# 工具函数
# ═══════════════════════════════════════════════════════════

def deep_merge(base, override):
    result = base.copy()
    for k, v in override.items():
        if k in result and isinstance(result[k], dict) and isinstance(v, dict):
            result[k] = deep_merge(result[k], v)
        else:
            result[k] = v
    return result

def load_config(config_path):
    cfg = DEFAULT_CONFIG.copy()
    if os.path.exists(config_path):
        try:
            with open(config_path, 'r', encoding='utf-8') as f:
                user_cfg = json.load(f)
            cfg = deep_merge(cfg, user_cfg)
            print(f"✓ 已加载外部配置：{config_path}")
        except Exception as e:
            print(f"⚠ 加载配置失败（{e}），使用默认配置")
    else:
        print("ℹ 未找到 config.json，使用内置默认配置")
    return cfg

def _find_cn_font():
    candidates = ["Microsoft YaHei", "SimHei", "SimSun", "KaiTi", "FangSong",
                  "Noto Sans CJK SC", "WenQuanYi Micro Hei", "PingFang SC"]
    available = {f.name for f in mpl_fm.fontManager.ttflist}
    for font in candidates:
        if font in available:
            return font
    return "sans-serif"

def sf(v):
    try:
        return float(v)
    except:
        return 0.0

def fm(v, d=0):
    try:
        return f"{float(v):,.{d}f}"
    except:
        return str(v) if v else "-"

def fp(v, d=2):
    try:
        f = float(v)
        return f"{f:.{d}f}%"
    except:
        return str(v) if v else "-"

# ═══════════════════════════════════════════════════════════
# 文档样式初始化（与 V8 一致）
# ═══════════════════════════════════════════════════════════

def setup_document_styles(doc, cfg_styles):
    font = cfg_styles["font_name"]
    body_sz = Pt(cfg_styles["body_size"])
    indent = Cm(cfg_styles.get("body_indent_cm", 0.74))

    def _set_font(style, size, bold=False, color=None, align=None, first_indent=None):
        style.font.name = font
        style.font.size = size
        style.font.bold = bold
        if color:
            style.font.color.rgb = RGBColor.from_string(color)
        style.element.rPr.rFonts.set(qn('w:eastAsia'), font)
        pf = style.paragraph_format
        if align is not None:
            pf.alignment = align
        if first_indent is not None:
            pf.first_line_indent = first_indent
        pf.space_after = Pt(4)

    n = doc.styles['Normal']
    _set_font(n, body_sz, first_indent=indent)
    _set_font(doc.styles['Heading 1'], Pt(cfg_styles["h1_size"]), bold=True)
    _set_font(doc.styles['Heading 2'], Pt(cfg_styles["h2_size"]), bold=True)
    _set_font(doc.styles['Heading 3'], Pt(cfg_styles["h2_size"]), bold=False)

    title_style = doc.styles.add_style('MiningTitle', WD_STYLE_TYPE.PARAGRAPH)
    _set_font(title_style, Pt(cfg_styles["title_size"]), bold=True,
              align=WD_ALIGN_PARAGRAPH.CENTER, first_indent=Cm(0))

    sub_style = doc.styles.add_style('MiningSubtitle', WD_STYLE_TYPE.PARAGRAPH)
    _set_font(sub_style, Pt(13), bold=False,
              align=WD_ALIGN_PARAGRAPH.CENTER, first_indent=Cm(0))

    ni_style = doc.styles.add_style('MiningNoIndent', WD_STYLE_TYPE.PARAGRAPH)
    _set_font(ni_style, Pt(9), first_indent=Cm(0))

    tc_style = doc.styles.add_style('MiningTableCell', WD_STYLE_TYPE.PARAGRAPH)
    _set_font(tc_style, Pt(cfg_styles["table_font_size"]), first_indent=Cm(0))

    th_style = doc.styles.add_style('MiningTableHeader', WD_STYLE_TYPE.PARAGRAPH)
    _set_font(th_style, Pt(cfg_styles["table_font_size"]), bold=True,
              color=cfg_styles["header_text_color"], first_indent=Cm(0))
    return cfg_styles

def add_para(doc, text, style='Normal'):
    return doc.add_paragraph(text, style=style)

def add_heading_styled(doc, text, level=1):
    return doc.add_heading(text, level=level)

def make_table(doc, headers, data_rows, cfg_styles):
    nrows = len(data_rows) + 1
    ncols = len(headers)
    tbl = doc.add_table(rows=nrows, cols=ncols)
    tbl.alignment = WD_TABLE_ALIGNMENT.CENTER
    tbl.autofit = True

    tblPr = tbl._tbl.tblPr
    borders = OE('w:tblBorders')
    for e in ['top', 'left', 'bottom', 'right', 'insideH', 'insideV']:
        b = OE(f'w:{e}')
        for a, v in [('val', 'single'), ('sz', '4'), ('color', '808080'), ('space', '0')]:
            b.set(qn(f'w:{a}'), v)
        borders.append(b)
    tblPr.append(borders)

    header_color = cfg_styles["header_bg_color"]
    for j, h in enumerate(headers):
        c = tbl.cell(0, j)
        c.text = ""
        p = c.paragraphs[0]
        p.style = doc.styles['MiningTableHeader']
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        r = p.add_run(str(h))
        shd = OE('w:shd')
        shd.set(qn('w:fill'), header_color)
        shd.set(qn('w:val'), 'clear')
        c._tc.get_or_add_tcPr().append(shd)

    for i, row in enumerate(data_rows):
        for j, val in enumerate(row):
            c = tbl.cell(i + 1, j)
            c.text = ""
            p = c.paragraphs[0]
            p.style = doc.styles['MiningTableCell']
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            r = p.add_run(str(val) if val is not None else "")
            if row[0] == "合计":
                r.bold = True
    return tbl

# ═══════════════════════════════════════════════════════════
# 列匹配器（不变）
# ═══════════════════════════════════════════════════════════

def auto_sort_steps(steps):
    """自动台阶排序：按数字前缀从小到大排列（如 1010, 1020, 1030...）"""
    def sort_key(s):
        m = re.match(r"(\d+)", str(s))
        return int(m.group(1)) if m else 99999
    return sorted(steps, key=sort_key)


class ColumnMatcher:
    def __init__(self, column_patterns):
        self.patterns = column_patterns

    def match(self, record, field_name):
        if field_name not in self.patterns:
            return None, None
        keyword_sets = self.patterns[field_name]
        if keyword_sets and isinstance(keyword_sets[0], str):
            keyword_sets = [keyword_sets]
        for k, v in record.items():
            if v is None:
                continue
            k_str = str(k)
            if all(any(kw in k_str for kw in kw_set) for kw_set in keyword_sets):
                return k_str, v
        return None, None

    def get(self, record, field_name):
        _, val = self.match(record, field_name)
        return val

    def get_or(self, record, field_name, default=None):
        val = self.get(record, field_name)
        return val if val is not None else default

    def match_header(self, headers, field_name):
        if field_name not in self.patterns:
            return None
        keyword_sets = self.patterns[field_name]
        if keyword_sets and isinstance(keyword_sets[0], str):
            keyword_sets = [keyword_sets]
        for idx, h in enumerate(headers):
            if h is None:
                continue
            h_str = str(h)
            if all(any(kw in h_str for kw in kw_set) for kw_set in keyword_sets):
                return idx
        return None

    def suggest_headers(self, headers, field_name):
        if field_name not in self.patterns:
            return []
        keyword_sets = self.patterns[field_name]
        if keyword_sets and isinstance(keyword_sets[0], str):
            keyword_sets = [keyword_sets]
        suggestions = []
        for h in headers:
            if h is None:
                continue
            h_str = str(h)
            score = sum(1 for kw_set in keyword_sets for kw in kw_set if kw in h_str)
            if score > 0:
                suggestions.append((h_str, score))
        suggestions.sort(key=lambda x: x[1], reverse=True)
        return [s[0] for s in suggestions[:5]]

# ═══════════════════════════════════════════════════════════
# V12 优化：公共表头解析逻辑（新增）
# ═══════════════════════════════════════════════════════════

def _find_header_and_parse(ws: Worksheet) -> Tuple[Optional[int], List[str]]:
    """
    V12 优化：在 worksheet 中查找表头行并解析 headers。
    返回 (h_idx, headers)，如果未找到表头行则返回 (None, [])。
    供 read_pinghua、read_duibai、classify_excel_file 等函数复用。
    """
    rows = list(ws.iter_rows(min_row=1, values_only=True))
    h_idx = None
    for i, r in enumerate(rows):
        if r and r[0] and is_header_row(r[0]):
            h_idx = i
            break
    if h_idx is None:
        return None, []
    headers = [str(h).strip() if h else "" for h in rows[h_idx]]
    return h_idx, headers

# ═══════════════════════════════════════════════════════════
# 数据读取（不变）
# ═══════════════════════════════════════════════════════════

def read_pinghua(filepath, config):
    cfg = config["columns"]["pinghua"]
    grouping = cfg.get("grouping", "every_n")
    n_per_step = cfg.get("n_per_step", 3)
    skip_kw = cfg.get("skip_keywords", [])
    col_match = ColumnMatcher(cfg.get("patterns", {}))
    wb = load_workbook(filepath, data_only=True)
    result = {}
    for sn in wb.sheetnames:
        ws = wb[sn]
        rows = list(ws.iter_rows(min_row=1, values_only=True))
        h_idx = None
        for i, r in enumerate(rows):
            if r and r[0] and is_header_row(r[0]):
                h_idx = i
                break
        if h_idx is None:
            continue
        headers = [str(h).strip() if h else "" for h in rows[h_idx]]
        data_rows = rows[h_idx + 1:]
        step_data = {}
        if grouping == "every_n":
            i = 0
            while i < len(data_rows):
                r0 = data_rows[i]
                if not r0 or r0[0] is None:
                    i += 1
                    continue
                nm = str(r0[0]).strip()
                if any(k in nm for k in skip_kw if k):
                    break
                rec = {}
                for offset in range(n_per_step):
                    if i + offset < len(data_rows) and data_rows[i + offset]:
                        row_vals = list(data_rows[i + offset])
                        for j, h in enumerate(headers):
                            if j < len(row_vals) and h:
                                key = f"{h}_r{offset}" if offset > 0 else h
                                rec[key] = row_vals[j]
                if n_per_step >= 3 and i + 2 < len(data_rows) and data_rows[i + 2]:
                    rec["品位"] = data_rows[i + 2][1] if len(data_rows[i + 2]) > 1 else None
                step_data[nm] = rec
                i += n_per_step
        result[sn] = step_data
    wb.close()
    return result

def read_duibai(filepath, config):
    cfg = config["columns"]["duibai"]
    sheet_patterns = cfg.get("sheet_patterns", [])
    summary_kw = cfg.get("summary_keywords", [])
    wb = load_workbook(filepath, data_only=True)
    result = {}
    for sn in wb.sheetnames:
        if not any(pat in sn for pat in sheet_patterns):
            continue
        ws = wb[sn]
        rows = list(ws.iter_rows(min_row=1, values_only=True))
        if not rows:
            continue
        r0, r1 = rows[0] if rows else [], rows[1] if len(rows) > 1 else []
        two_level = (r0 and r0[0] and is_header_row(r0[0]) and
                     r1 and r1[0] is None and any(c is not None for c in r1))
        if two_level:
            parent = [str(c).strip() if c else None for c in r0]
            child = [str(c).strip() if c else "" for c in r1]
            last = ""
            for i in range(len(parent)):
                parent[i] = parent[i] if parent[i] else last
                if parent[i]:
                    last = parent[i]
            headers = [f"{p}_{c}" if p and c else (p or c) for p, c in zip(parent, child)]
            seen = {}
            final = []
            for h in headers:
                if h in seen:
                    seen[h] += 1
                    final.append(f"{h}_{seen[h]}")
                else:
                    seen[h] = 1
                    final.append(h)
            headers = final
            data_start = 2
        else:
            data_start = None
            for i, r in enumerate(rows):
                if r and r[0] and is_header_row(r[0]):
                    data_start = i
                    break
            if data_start is None:
                continue
            headers = [str(c).strip() if c else "" for c in rows[data_start]]
            data_start += 1
        records = []
        for row in rows[data_start:]:
            if not row or row[0] is None:
                continue
            nm = str(row[0]).strip()
            if not nm:
                continue
            rec = {"_sheet": sn, "_step": nm}
            for j, h in enumerate(headers):
                if h:
                    rec[h] = row[j] if j < len(row) else None
                else:
                    rec[f"col_{j}"] = row[j] if j < len(row) else None
            rec["_is_summary"] = any(k in nm for k in summary_kw)
            records.append(rec)
        result[sn] = records
    wb.close()
    return result

def is_step_row(rec):
    return not rec.get("_is_summary", False)

# V12 优化：validate_ore_logic 循环外预创建 ColumnMatcher，循环内复用
def validate_ore_logic(ore_all, col_match, label=""):
    fixed = 0
    # V12 优化：在循环外预创建两个 ColumnMatcher，避免循环内反复创建
    cm_prod = ColumnMatcher({"生产_金属量": [["生产"], ["金属量"]]})
    cm_zhtk = ColumnMatcher({"ZHTK_金属量": [["ZHTK"], ["金属量"]]})
    for rec in ore_all:
        q = sf(col_match.get_or(rec, "生产_矿石量"))
        c = sf(col_match.get_or(rec, "生产_品位"))
        p = sf(col_match.get_or(rec, "生产_金属量"))
        z_q = sf(col_match.get_or(rec, "ZHTK_矿石量"))
        z_c = sf(col_match.get_or(rec, "ZHTK_品位"))
        z_p = sf(col_match.get_or(rec, "ZHTK_金属量"))

        for q_val, c_val, p_val, key_field, cm in [
            (q, c, p, "生产_金属量", cm_prod), (z_q, z_c, z_p, "ZHTK_金属量", cm_zhtk)
        ]:
            if not q_val or not c_val or not p_val:
                continue
            expected = q_val * c_val / 100.0
            if abs(expected - p_val) > 0.5:
                for k in rec:
                    kk, _ = cm.match({k: rec[k]}, key_field)
                    if kk is not None:
                        rec[k] = expected
                        fixed += 1
                        break
    return fixed

# ═══════════════════════════════════════════════════════════
# ═══════════════════════════════════════════════════════════
# 🎨 图表工厂层（V9.1 真正解耦）
# 所有 make_*_chart 函数统一返回 (fig, buf)
# ═══════════════════════════════════════════════════════════

_FONT = _find_cn_font()
plt.rcParams.update({
    'font.family': 'sans-serif',
    'font.sans-serif': [_FONT, 'DejaVu Sans'],
    'axes.unicode_minus': False,
    'figure.dpi': 100,
    'savefig.dpi': 200,
})

def _fig_to_buf(fig, dpi=200):
    buf = BytesIO()
    fig.savefig(buf, format='png', dpi=dpi, bbox_inches='tight', facecolor='white')
    buf.seek(0)
    return buf

def _style_ax(ax, title="", xlabel="", ylabel="", ylim=None, grid_y=True, grid_x=False):
    ax.set_title(title, fontsize=14, fontweight='normal', pad=12)
    ax.set_xlabel(xlabel, fontsize=10)
    ax.set_ylabel(ylabel, fontsize=10)
    if ylim:
        ax.set_ylim(ylim)
    for s in ax.spines.values():
        s.set_color("#555555")
        s.set_linewidth(0.8)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.tick_params(colors="#555555", labelsize=9)
    ax.tick_params(axis='x', rotation=0, pad=6)   # 强制水平
    if grid_y:
        ax.grid(True, axis='y', linestyle='--', color="#E2E6EA", alpha=0.3, linewidth=0.6)
    if grid_x:
        ax.grid(True, axis='x', linestyle='--', color="#E2E6EA", alpha=0.3, linewidth=0.6)
    ax.set_axisbelow(True)

def _add_bar_labels(ax, bars, fmt="{:.0f}", bar_color=None, fontsize=7, pad=0.25):
    for b in bars:
        h = b.get_height()
        if h == 0:
            continue
        color = bar_color or b.get_facecolor()
        y_range = ax.get_ylim()[1] - ax.get_ylim()[0]
        ax.text(
            b.get_x() + b.get_width() / 2,
            h + y_range * 0.015,
            fmt.format(h),
            ha='center', va='bottom', fontsize=fontsize,
            color='#FFFFFF', fontweight='bold',
            bbox=dict(boxstyle=f'round,pad={pad}', facecolor=color, edgecolor='none', alpha=0.92)
        )

def _add_line_labels(ax, x_vals, y_vals, color, fontsize=6.5, direction='up', ylim=None, fmt="{:.3f}"):
    if ylim is None:
        ylim = ax.get_ylim()
    y_span = ylim[1] - ylim[0]
    n = len(x_vals)

    if direction is None:
        # 智能波峰/波谷检测
        for i, (x, y) in enumerate(zip(x_vals, y_vals)):
            if y is None:
                continue
            prev_y = y_vals[i-1] if i > 0 else y
            next_y = y_vals[i+1] if i < n-1 else y
            if y >= prev_y and y >= next_y:
                va = 'bottom'
                offset = y_span * 0.04
            elif y <= prev_y and y <= next_y:
                va = 'top'
                offset = -y_span * 0.04
            else:
                va = 'bottom'
                offset = y_span * 0.04
            ax.text(x, y + offset, fmt.format(y),
                    ha='center', va=va, fontsize=fontsize,
                    color='#FFFFFF', fontweight='bold',
                    bbox=dict(boxstyle='round,pad=0.2', facecolor=color, edgecolor='none', alpha=0.88))
    else:
        offset_ratio = 0.04
        for x, y in zip(x_vals, y_vals):
            if y is None:
                continue
            if direction == 'up':
                vy = y + y_span * offset_ratio
                va = 'bottom'
            else:
                vy = y - y_span * offset_ratio
                va = 'top'
            ax.text(x, vy, fmt.format(y),
                    ha='center', va=va, fontsize=fontsize,
                    color='#FFFFFF', fontweight='bold',
                    bbox=dict(boxstyle='round,pad=0.2', facecolor=color, edgecolor='none', alpha=0.88))

# ═══════════════════════════════════════════════════════════
# 5. 通用图表工厂（V10 核心：2 个通用函数驱动 8 张图表）
# ═══════════════════════════════════════════════════════════

def draw_comparison_bars(step_names, data1, data2, label1, label2,
                         title, ylabel, color1, color2, fmt_key="ore",
                         unit_scale=1.0, bar_label_fmt=None) -> Tuple[Figure, BytesIO]:
    """通用双模型对比柱状图：自动单位换算、标签、动态图例"""
    if bar_label_fmt is None:
        bar_label_fmt = {"ore": "{:.1f}", "metal": "{:.2f}", "area": "{:.0f}"}
    fmt = bar_label_fmt.get(fmt_key, "{:.1f}")
    n = len(step_names)
    if n == 0:
        warnings.warn("数据为空，无法绘制图表")
        fig, ax = plt.subplots()
        ax.text(0.5, 0.5, "无数据", ha='center', va='center')
        return fig, _fig_to_buf(fig)

    fig, ax = plt.subplots(figsize=(12, 5.5))
    x = range(n)
    w = 0.35
    d1 = [v / unit_scale for v in data1]
    d2 = [v / unit_scale for v in data2]

    bars1 = ax.bar([i - w/2 for i in x], d1, w, color=color1, label=label1, edgecolor='white', linewidth=0.5)
    bars2 = ax.bar([i + w/2 for i in x], d2, w, color=color2, label=label2, edgecolor='white', linewidth=0.5)

    ax.set_xticks(x)
    ax.set_xticklabels(step_names, rotation=0, ha='center')
    _style_ax(ax, title, '台阶', ylabel)
    _add_bar_labels(ax, bars1, fmt=fmt, bar_color=color1, pad=0.15, fontsize=6)
    _add_bar_labels(ax, bars2, fmt=fmt, bar_color=color2, pad=0.15, fontsize=6)

    if len(step_names) <= 6:
        ax.legend(loc='upper right', frameon=True, facecolor='white', edgecolor='#cccccc', fontsize=8)
    else:
        ax.legend(loc='upper left', bbox_to_anchor=(0.01, 0.99), frameon=True, facecolor='white', edgecolor='#cccccc', fontsize=8)
    plt.tight_layout()
    return fig, _fig_to_buf(fig)


def draw_comparison_lines(step_names, data1, data2, label1, label2,
                          title, ylabel, color1, color2) -> Tuple[Figure, BytesIO]:
    """通用双模型对比折线图：智能标签防遮挡"""
    n = len(step_names)
    if n == 0:
        warnings.warn("数据为空")
        fig, ax = plt.subplots()
        ax.text(0.5, 0.5, "无数据", ha='center', va='center')
        return fig, _fig_to_buf(fig)

    fig, ax = plt.subplots(figsize=(12, 5))
    x = range(n)
    ax.plot(x, data1, 'o-', color=color1, linewidth=2.2, markersize=8, label=label1)
    ax.plot(x, data2, 's--', color=color2, linewidth=2.2, markersize=8, label=label2)
    _add_line_labels(ax, x, data1, color1, direction=None)
    _add_line_labels(ax, x, data2, color2, direction=None)
    ax.set_xticks(x)
    ax.set_xticklabels(step_names, rotation=0, ha='center')
    _style_ax(ax, title, '台阶', ylabel, grid_x=True)

    if len(step_names) <= 6:
        ax.legend(loc='upper right', frameon=True, facecolor='white', edgecolor='#cccccc', fontsize=8)
    else:
        ax.legend(loc='upper left', bbox_to_anchor=(0.01, 0.99), frameon=True, facecolor='white', edgecolor='#cccccc', fontsize=8)
    plt.tight_layout()
    return fig, _fig_to_buf(fig)


# ===== 面积组合图（保留 V9.1 完整双Y轴版） =====
def make_area_dr_chart(step_names, area_data, cfg):
    """矿体面积及重合面积组合图 — 3柱+2线，双Y轴，带数据校验、动态图例、动态刻度"""
    colors = cfg.get("colors", {})
    bar_fmts = cfg.get("bar_label_fmt", {})
    n = len(step_names)
    if n == 0:
        warnings.warn("面积数据为空，无法绘制图表")
        fig, ax = plt.subplots()
        ax.text(0.5, 0.5, "无数据", ha='center', va='center')
        return fig, _fig_to_buf(fig)

    prod_vals = [area_data.get(s, {}).get("prod", 0) for s in step_names]
    zhtk_vals = [area_data.get(s, {}).get("zhtk", 0) for s in step_names]
    over_vals = [area_data.get(s, {}).get("over", 0) for s in step_names]
    rel_e_vals = [area_data.get(s, {}).get("rel_e", 0) for s in step_names]
    dr_vals = [area_data.get(s, {}).get("dr", 0) for s in step_names]

    # 数据完整性校验
    if not any(prod_vals) and not any(zhtk_vals):
        warnings.warn("面积数据全为0，请检查数据源")

    fig, ax1 = plt.subplots(figsize=(14, 5.8))
    ax2 = ax1.twinx()
    x = range(n)
    w = 0.22

    ax1.bar([i - w for i in x], prod_vals, w, color=colors.get("area_prod", "#3B76B8"),
            label='生产模型', edgecolor='white', linewidth=0.5)
    ax1.bar(x, zhtk_vals, w, color=colors.get("area_zhtk", "#FF901C"),
            label='ZHTK模型', edgecolor='white', linewidth=0.5)
    ax1.bar([i + w for i in x], over_vals, w, color=colors.get("area_overlap", "#FFC800"),
            label='重合面积', edgecolor='white', linewidth=0.5)

    ax2.plot(x, rel_e_vals, 'o-', color=colors.get("area_rel_e", "#235E1F"), linewidth=2,
             markersize=6, label='相对误差', zorder=5)
    ax2.plot(x, dr_vals, 's-', color=colors.get("area_dr", "#649DC9"), linewidth=2,
             markersize=6, label='面积重合率', zorder=5)

    # 动态智能标签（direction=None启用波峰波谷检测）
    _add_line_labels(ax2, x, rel_e_vals, colors.get("area_rel_e", "#235E1F"), direction=None, fmt="{:.0f}")
    _add_line_labels(ax2, x, dr_vals, colors.get("area_dr", "#649DC9"), direction=None, fmt="{:.0f}")

    ax1.set_xticks(x)
    ax1.set_xticklabels(step_names, fontsize=9)
    ax1.tick_params(axis='y', colors='#333333', labelsize=9)
    ax2.tick_params(axis='y', colors='#333333', labelsize=9)

    # ===== 动态次轴刻度（V9改进） =====
    all_vals = rel_e_vals + dr_vals
    if all_vals:
        min_val = min(all_vals)
        max_val = max(all_vals)
        ymin = min_val - abs(min_val) * 0.3 if min_val < 0 else min_val * 0.7
        ymax = max_val + abs(max_val) * 0.3
        # 取整十
        ymin = int(ymin / 10) * 10
        ymax = int(ymax / 10) * 10 + (10 if ymax % 10 else 0)
        ax2.set_ylim(ymin, ymax)
        ax2.yaxis.set_major_locator(AutoLocator())
    else:
        ax2.set_ylim(-100, 100)

    ax1.grid(axis='y', color='#D9D9D9', linestyle='--', linewidth=0.6, alpha=0.8)
    ax1.set_axisbelow(True)

    title = f'{cfg.get("mine", {}).get("name", "")} 各台阶矿体面积及重合面积图'
    ax1.set_title(title, fontsize=12, fontweight='normal', pad=8)
    ax1.set_ylabel('面积 (m²)', fontsize=9, color='#333333')
    ax1.ticklabel_format(axis='y', style='scientific', scilimits=(0, 0))
    ax1.yaxis.get_offset_text().set_fontsize(8)

    lines1, labels1 = ax1.get_legend_handles_labels()
    lines2, labels2 = ax2.get_legend_handles_labels()
    # ===== 动态图例（V9改进） =====
    if len(step_names) <= 6:
        ax1.legend(lines1 + lines2, labels1 + labels2,
                   loc='upper right', frameon=True, facecolor='white', edgecolor='#CCCCCC', fontsize=8)
    else:
        ax1.legend(lines1 + lines2, labels1 + labels2,
                   loc='upper center', bbox_to_anchor=(0.5, -0.08),
                   ncol=5, frameon=True, facecolor='white', edgecolor='#CCCCCC', fontsize=8)
    plt.tight_layout()
    return fig, _fig_to_buf(fig)

def make_area_wr_chart(step_names, area_data, cfg):
    """歪曲面积及歪曲率组合图 — 堆叠柱+折线，独立配色"""
    colors = cfg.get("colors", {})
    n = len(step_names)
    if n == 0:
        warnings.warn("歪曲数据为空")
        fig, ax = plt.subplots()
        ax.text(0.5, 0.5, "无数据", ha='center', va='center')
        return fig, _fig_to_buf(fig)

    wadd_vals = [area_data.get(s, {}).get("wadd", 0) for s in step_names]
    wsub_vals = [area_data.get(s, {}).get("wsub", 0) for s in step_names]
    wr_vals = [area_data.get(s, {}).get("wr", 0) for s in step_names]

    if not any(wadd_vals) and not any(wsub_vals):
        warnings.warn("歪曲面积数据全为0")

    fig, ax1 = plt.subplots(figsize=(14, 5.8))
    ax2 = ax1.twinx()
    x = range(n)

    ax1.bar(x, wadd_vals, 0.55, color=colors.get("area_wadd", "#649DC9"),
            label='增加', edgecolor='white', linewidth=0.5)
    # ===== V9改进：减少柱使用独立配色 =====
    ax1.bar(x, [-v for v in wsub_vals], 0.55, color=colors.get("area_decrease", "#A0A0A0"),
            bottom=0, label='减少', edgecolor='white', linewidth=0.5)

    ax2.plot(x, wr_vals, 'o-', color=colors.get("area_wr", "#2F4F2E"), linewidth=2,
             markersize=6, label='歪曲率', zorder=5)

    _add_line_labels(ax2, x, wr_vals, colors.get("area_wr", "#2F4F2E"), direction=None, fmt="{:.0f}")

    ax1.set_xticks(x)
    ax1.set_xticklabels(step_names, fontsize=9)
    ax1.tick_params(axis='y', colors='#333333', labelsize=9)
    ax2.tick_params(axis='y', colors='#333333', labelsize=9)

    # 动态次轴刻度
    if wr_vals:
        max_wr = max(wr_vals)
        ymax = max_wr + abs(max_wr) * 0.3
        ymax = int(ymax / 10) * 10 + (10 if ymax % 10 else 0)
        ax2.set_ylim(-ymax, ymax)
        ax2.yaxis.set_major_locator(AutoLocator())
    else:
        ax2.set_ylim(-70, 70)
    ax2.set_ylabel('歪曲率 (%)', fontsize=9, color='#333333')

    ax1.axhline(y=0, color='#999999', linewidth=0.8, linestyle='-', alpha=0.6)
    ax1.grid(axis='y', color='#D9D9D9', linestyle='--', linewidth=0.6, alpha=0.8)
    ax1.set_axisbelow(True)

    title = f'{cfg.get("mine", {}).get("name", "")} 各台阶歪曲面积及歪曲率'
    ax1.set_title(title, fontsize=12, fontweight='normal', pad=8)
    ax1.set_ylabel('面积差值 (m²)', fontsize=9, color='#333333')

    lines1, labels1 = ax1.get_legend_handles_labels()
    lines2, labels2 = ax2.get_legend_handles_labels()
    # 动态图例
    if len(step_names) <= 6:
        ax1.legend(lines1 + lines2, labels1 + labels2, loc='upper right', frameon=True,
                   facecolor='white', edgecolor='#CCCCCC', fontsize=8)
    else:
        ax1.legend(lines1 + lines2, labels1 + labels2,
                   loc='upper center', bbox_to_anchor=(0.5, -0.08),
                   ncol=3, frameon=True, facecolor='white', edgecolor='#CCCCCC', fontsize=8)
    plt.tight_layout()
    return fig, _fig_to_buf(fig)

# ═══════════════════════════════════════════════════════════
# V12 优化：统一图表函数（合并 6 个重复图表函数为 3 个）
# ═══════════════════════════════════════════════════════════

def make_ore_chart(step_names, step_detail, cfg, variant="detail") -> Tuple[Figure, BytesIO]:
    """
    V12 优化：统一矿石量对比柱状图。
    variant="detail" 时标签为 "生产勘探模型矿石量 (万t)"（原 make_ore_qty_chart），
    variant="step" 时为 "生产勘探模型 (万t)"（原 make_step_ore_chart）。
    """
    colors = cfg.get("colors", {})
    bar_fmts = cfg.get("bar_label_fmt", {})
    if variant == "step":
        label1 = "生产勘探模型 (万t)"
        label2 = "ZHTK地质模型 (万t)"
    else:
        label1 = "生产勘探模型矿石量 (万t)"
        label2 = "ZHTK地质模型矿石量 (万t)"
    return draw_comparison_bars(
        step_names,
        [step_detail.get(s, {}).get("p_q", 0) for s in step_names],
        [step_detail.get(s, {}).get("z_q", 0) for s in step_names],
        label1, label2,
        f'{cfg.get("mine", {}).get("name", "")} 各台阶矿石量', "矿石量 (万t)",
        colors.get("model_prod", "#2E6B8A"), colors.get("model_zhtk", "#C98247"),
        "ore", unit_scale=10000, bar_label_fmt=bar_fmts
    )

def make_grade_chart(step_names, step_detail, cfg, variant="detail") -> Tuple[Figure, BytesIO]:
    """
    V12 优化：统一品位对比折线图。
    variant="detail" 时标签为 "生产勘探模型品位 (%)"（原 make_grade_line_chart），
    variant="step" 时为 "生产勘探模型品位 (%)"（原 make_step_grade_chart）。
    注：两者标签实际相同，保留 variant 参数以备扩展。
    """
    colors = cfg.get("colors", {})
    return draw_comparison_lines(
        step_names,
        [step_detail.get(s, {}).get("p_c", 0) for s in step_names],
        [step_detail.get(s, {}).get("z_c", 0) for s in step_names],
        "生产勘探模型品位 (%)", "ZHTK地质模型品位 (%)",
        f'{cfg.get("mine", {}).get("name", "")} 各台阶品位', "品位 (%)",
        colors.get("model_prod", "#2E6B8A"), colors.get("model_zhtk", "#C98247")
    )

def make_metal_chart(step_names, step_detail, cfg, variant="detail") -> Tuple[Figure, BytesIO]:
    """
    V12 优化：统一金属量对比柱状图。
    variant="detail" 时标签为 "生产勘探模型金属量 (万t)"（原 make_metal_chart），
    variant="step" 时为 "生产勘探模型 (万t)"（原 make_step_metal_chart）。
    """
    colors = cfg.get("colors", {})
    bar_fmts = cfg.get("bar_label_fmt", {})
    if variant == "step":
        label1 = "生产勘探模型 (万t)"
        label2 = "ZHTK地质模型 (万t)"
    else:
        label1 = "生产勘探模型金属量 (万t)"
        label2 = "ZHTK地质模型金属量 (万t)"
    return draw_comparison_bars(
        step_names,
        [step_detail.get(s, {}).get("p_p", 0) for s in step_names],
        [step_detail.get(s, {}).get("z_p", 0) for s in step_names],
        label1, label2,
        f'{cfg.get("mine", {}).get("name", "")} 各台阶金属量', "金属量 (万t)",
        colors.get("model_prod", "#2E6B8A"), colors.get("model_zhtk", "#C98247"),
        "metal", unit_scale=10000, bar_label_fmt=bar_fmts
    )

# ═══════════════════════════════════════════════════════════
# V12 优化：兼容 wrapper（保留旧函数名，内部调用统一函数）
# ═══════════════════════════════════════════════════════════

def make_ore_qty_chart(step_names, step_detail, cfg):
    """矿石量对比柱状图（向后兼容 wrapper → make_ore_chart variant='detail'）"""
    return make_ore_chart(step_names, step_detail, cfg, variant="detail")

def make_grade_line_chart(step_names, step_detail, cfg):
    """品位对比折线图（向后兼容 wrapper → make_grade_chart variant='detail'）"""
    return make_grade_chart(step_names, step_detail, cfg, variant="detail")

def make_step_ore_chart(step_names, step_detail, cfg):
    """各台阶矿石量柱状图（向后兼容 wrapper → make_ore_chart variant='step'）"""
    return make_ore_chart(step_names, step_detail, cfg, variant="step")

def make_step_grade_chart(step_names, step_detail, cfg):
    """各台阶品位折线图（向后兼容 wrapper → make_grade_chart variant='step'）"""
    return make_grade_chart(step_names, step_detail, cfg, variant="step")

def make_step_metal_chart(step_names, step_detail, cfg):
    """各台阶金属量柱状图（向后兼容 wrapper → make_metal_chart variant='step'）"""
    return make_metal_chart(step_names, step_detail, cfg, variant="step")

# ═══════════════════════════════════════════════════════════
# 3. 智能识别与数据提取（V10 核心）
# ═══════════════════════════════════════════════════════════

def build_synonym_regex():
    patterns = {}
    for std_name, synonyms in SYNONYM_MAP.items():
        regex = '|'.join(re.escape(s) for s in synonyms)
        patterns[std_name] = re.compile(regex, re.IGNORECASE)
    return patterns

SYNONYM_REGEX = build_synonym_regex()

def score_headers(headers, file_type="duibai"):
    required_fields = []
    if file_type == "duibai":
        required_fields = [
            "生产_矿石量", "ZHTK_矿石量", "生产_品位", "ZHTK_品位",
            "生产_金属量", "ZHTK_金属量", "面积_生产", "面积_ZHTK",
            "面积_重合", "面积_歪曲增加", "面积_歪曲减少"
        ]
    elif file_type == "pinghua":
        required_fields = ["贫化率", "损失率", "出矿量", "品位"]

    field_col_map = {}
    for field in required_fields:
        best_score = 0
        best_idx = None
        regex = SYNONYM_REGEX.get(field)
        if not regex:
            continue
        for idx, h in enumerate(headers):
            if h is None:
                continue
            h_str = str(h).strip()
            if not h_str:
                continue
            if regex.search(h_str):
                score = len(regex.findall(h_str)) * 2 + 1.0 / (len(h_str) + 1)
                if score > best_score:
                    best_score = score
                    best_idx = idx
        if best_score > 0:
            field_col_map[field] = best_idx
    return field_col_map

def classify_excel_file(filepath: str) -> dict:
    wb = load_workbook(filepath, data_only=True)
    result = {"type": None, "sheet_mappings": {}, "all_headers": {}}
    for sn in wb.sheetnames:
        ws = wb[sn]
        # V12 优化：复用 _find_header_and_parse
        h_idx, headers = _find_header_and_parse(ws)
        if h_idx is None:
            continue
        result["all_headers"][sn] = headers

    type_scores = {"duibai": 0, "pinghua": 0}
    for sn, headers in result["all_headers"].items():
        for ftype in ["duibai", "pinghua"]:
            mapping = score_headers(headers, ftype)
            total = 11 if ftype == "duibai" else 4
            coverage = len(mapping) / total if total > 0 else 0
            type_scores[ftype] = max(type_scores[ftype], coverage)

    best_type = None
    best_score = 0
    for ftype, score in type_scores.items():
        if score > FILE_SCORE_THRESHOLD and score > best_score:
            best_score = score
            best_type = ftype

    if best_type:
        result["type"] = best_type
        for sn, headers in result["all_headers"].items():
            result["sheet_mappings"][sn] = score_headers(headers, best_type)

    wb.close()
    return result

# V12 优化：新增 classify_excel_file_cached，接受可选 wb 参数避免重复打开
def classify_excel_file_cached(filepath: str, wb=None) -> dict:
    """
    V12 优化：classify_excel_file 的缓存友好版本。
    若传入 wb（已打开的 workbook），则不再重新打开文件。
    返回格式与 classify_excel_file 完全一致。
    """
    should_close = False
    if wb is None:
        wb = load_workbook(filepath, data_only=True)
        should_close = True

    result = {"type": None, "sheet_mappings": {}, "all_headers": {}}
    for sn in wb.sheetnames:
        ws = wb[sn]
        h_idx, headers = _find_header_and_parse(ws)
        if h_idx is None:
            continue
        result["all_headers"][sn] = headers

    type_scores = {"duibai": 0, "pinghua": 0}
    for sn, headers in result["all_headers"].items():
        for ftype in ["duibai", "pinghua"]:
            mapping = score_headers(headers, ftype)
            total = 11 if ftype == "duibai" else 4
            coverage = len(mapping) / total if total > 0 else 0
            type_scores[ftype] = max(type_scores[ftype], coverage)

    best_type = None
    best_score = 0
    for ftype, score in type_scores.items():
        if score > FILE_SCORE_THRESHOLD and score > best_score:
            best_score = score
            best_type = ftype

    if best_type:
        result["type"] = best_type
        for sn, headers in result["all_headers"].items():
            result["sheet_mappings"][sn] = score_headers(headers, best_type)

    if should_close:
        wb.close()
    return result

# V12 优化：extract_duibai_data 内部只打开一次 workbook
def extract_duibai_data(filepath: str, cfg: dict) -> Dict[str, List[dict]]:
    # V12 优化：只打开一次 workbook，传给 classify_excel_file_cached 复用
    wb = load_workbook(filepath, data_only=True)
    info = classify_excel_file_cached(filepath, wb=wb)
    if info["type"] != "duibai":
        wb.close()
        raise ValueError(f"文件 {os.path.basename(filepath)} 未被识别为探采对比表")

    db_data = {}
    summary_kw = cfg["columns"]["duibai"].get("summary_keywords", ["合计"])

    for sn, mapping in info["sheet_mappings"].items():
        if not mapping:
            continue
        ws = wb[sn]
        rows = list(ws.iter_rows(min_row=1, values_only=True))
        h_idx = None
        for i, r in enumerate(rows):
            if r and r[0] and is_header_row(r[0]):
                h_idx = i
                break
        if h_idx is None:
            continue
        data_rows = rows[h_idx+1:]
        records = []
        for row in data_rows:
            if not row or row[0] is None:
                continue
            nm = str(row[0]).strip()
            if not nm:
                continue
            rec = {"_sheet": sn, "_step": nm}
            for field, col_idx in mapping.items():
                rec[field] = row[col_idx] if col_idx < len(row) else None
            rec["_is_summary"] = any(k in nm for k in summary_kw)
            records.append(rec)
        db_data[sn] = records

    wb.close()
    if not db_data:
        raise ValueError("未能提取到有效数据")
    return db_data

# V12 优化：extract_pinghua_data 内部只打开一次 workbook
def extract_pinghua_data(filepath: str, cfg: dict) -> Dict[str, dict]:
    # V12 优化：只打开一次 workbook，传给 classify_excel_file_cached 复用
    wb = load_workbook(filepath, data_only=True)
    info = classify_excel_file_cached(filepath, wb=wb)
    if info["type"] != "pinghua":
        wb.close()
        raise ValueError(f"文件 {os.path.basename(filepath)} 未被识别为贫化损失表")

    ph_data = {}
    pinghua_cfg = cfg["columns"]["pinghua"]
    grouping = pinghua_cfg.get("grouping", "every_n")
    n_per_step = pinghua_cfg.get("n_per_step", 3)
    skip_kw = pinghua_cfg.get("skip_keywords", [])

    for sn, mapping in info["sheet_mappings"].items():
        if not mapping:
            continue
        ws = wb[sn]
        rows = list(ws.iter_rows(min_row=1, values_only=True))
        h_idx = None
        for i, r in enumerate(rows):
            if r and r[0] and is_header_row(r[0]):
                h_idx = i
                break
        if h_idx is None:
            continue
        data_rows = rows[h_idx+1:]
        step_data = {}

        if grouping == "every_n":
            i = 0
            while i < len(data_rows):
                r0 = data_rows[i]
                if not r0 or r0[0] is None:
                    i += 1
                    continue
                nm = str(r0[0]).strip()
                if any(k in nm for k in skip_kw if k):
                    break
                rec = {}
                for offset in range(n_per_step):
                    if i + offset < len(data_rows):
                        row_vals = data_rows[i + offset]
                        for field, col_idx in mapping.items():
                            if col_idx < len(row_vals):
                                key = f"{field}_{offset}" if offset > 0 else field
                                rec[key] = row_vals[col_idx]
                step_data[nm] = rec
                i += n_per_step

        ph_data[sn] = step_data

    wb.close()
    if not ph_data:
        raise ValueError("未能提取到有效数据")
    return ph_data

# ═══════════════════════════════════════════════════════════
# 4. SQLite 历史数据库（V10 新增）
# ═══════════════════════════════════════════════════════════

def init_database(db_path):
    conn = sqlite3.connect(db_path)
    cursor = conn.cursor()
    cursor.execute('''
        CREATE TABLE IF NOT EXISTS step_indicators (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            year INTEGER,
            quarter INTEGER,
            step TEXT,
            indicator_type TEXT,
            value REAL,
            UNIQUE(year, quarter, step, indicator_type)
        )
    ''')
    cursor.execute('''
        CREATE TABLE IF NOT EXISTS total_indicators (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            year INTEGER,
            quarter INTEGER,
            indicator_type TEXT,
            value REAL,
            UNIQUE(year, quarter, indicator_type)
        )
    ''')
    conn.commit()
    conn.close()

# V12 修复：save_indicators_to_db 不再跳过核心指标，改为排除非核心计算字段
def save_indicators_to_db(db_path: str, year: int, quarter: int,
                          step_data: dict, total_data: dict) -> None:
    conn = sqlite3.connect(db_path)
    cursor = conn.cursor()
    # V12 优化：排除非核心计算字段，而非跳过核心指标
    _exclude_fields = {"q_diff", "q_diff_pct"}
    for step, d in step_data.items():
        if step == "合计":
            continue

        for key, val in d.items():
            # V12 修复：只排除非核心计算字段，核心指标（p_q/z_q/p_c/z_c/p_p/z_p）正常入库
            if key in _exclude_fields:
                continue
            if isinstance(val, (int, float)):
                cursor.execute(
                    "INSERT OR REPLACE INTO step_indicators (year, quarter, step, indicator_type, value) VALUES (?, ?, ?, ?, ?)",
                    (year, quarter, step, key, val)
                )
    for key, val in total_data.items():
        if isinstance(val, (int, float)):
            cursor.execute(
                "INSERT OR REPLACE INTO total_indicators (year, quarter, indicator_type, value) VALUES (?, ?, ?, ?)",
                (year, quarter, key, val)
            )
    conn.commit()
    conn.close()

# V12 修复：detect_anomaly 合并为单次连接、单条 SQL
def detect_anomaly(db_path: str, year: int, quarter: int,
                   indicator_type: str, threshold: int = 30) -> Tuple[bool, float, Optional[float]]:
    # V12 优化：一次连接、一条 SQL 用 ORDER BY + LIMIT 2 获取最近两期数据
    conn = sqlite3.connect(db_path)
    cursor = conn.cursor()

    # 计算上一期的年份和季度
    if quarter > 1:
        prev_year, prev_quarter = year, quarter - 1
    else:
        prev_year, prev_quarter = year - 1, 4

    # 用一条 SQL 获取当前期和上一期的数据（按年份、季度降序排列，取最近2条）
    cursor.execute(
        """SELECT year, quarter, value FROM total_indicators
           WHERE indicator_type = ? AND ((year = ? AND quarter = ?) OR (year = ? AND quarter = ?))
           ORDER BY year DESC, quarter DESC
           LIMIT 2""",
        (indicator_type, year, quarter, prev_year, prev_quarter)
    )
    rows = cursor.fetchall()
    conn.close()

    if len(rows) < 2:
        if len(rows) == 1:
            # 只有当前期，没有上一期
            return False, 0, None
        return False, 0, None

    # rows[0] 是当前期（最新），rows[1] 是上一期
    curr_val = rows[0][2]
    prev_val = rows[1][2]

    if prev_val == 0:
        return False, 0, prev_val
    change = (curr_val - prev_val) / prev_val * 100
    if abs(change) > threshold:
        return True, change, prev_val
    return False, change, prev_val

# ═══════════════════════════════════════════════════════════
# 预检函数（不变）
# ═══════════════════════════════════════════════════════════

def verify_excel_structure(f_pinghua, f_duibai, cfg):
    errors = []
    ph_cfg = cfg["columns"]["pinghua"]
    db_cfg = cfg["columns"]["duibai"]

    if not os.path.exists(f_pinghua):
        errors.append(f"找不到贫化损失统计表：{f_pinghua}")
    else:
        try:
            wb = load_workbook(f_pinghua, data_only=True)
            sheets = wb.sheetnames
            required = ph_cfg.get("required_sheets")
            if required:
                for rs in required:
                    if rs not in sheets:
                        errors.append(f"贫化损失表中缺少「{rs}」工作表（现有：{', '.join(sheets)}）")
            ph_cm = ColumnMatcher(ph_cfg.get("patterns", {}))
            required_cols = ph_cfg.get("required", [])
            for sn in sheets:
                ws = wb[sn]
                rows = list(ws.iter_rows(min_row=1, values_only=True))
                h_idx = None
                for i, r in enumerate(rows):
                    if r and r[0] and is_header_row(r[0]):
                        h_idx = i
                        break
                if h_idx is None:
                    continue
                headers = [str(h).strip() if h else "" for h in rows[h_idx]]
                for rc in required_cols:
                    if ph_cm.match_header(headers, rc) is None:
                        suggestions = ph_cm.suggest_headers(headers, rc)
                        sug_str = f"（建议检查是否包含：{', '.join(suggestions)}）" if suggestions else ""
                        errors.append(f"贫化损失表「{sn}」缺少「{rc}」列{sug_str}")
            wb.close()
        except Exception as e:
            errors.append(f"读取贫化损失表失败：{e}")

    if not os.path.exists(f_duibai):
        errors.append(f"找不到探采对比表格：{f_duibai}")
    else:
        try:
            wb = load_workbook(f_duibai, data_only=True)
            sheets = wb.sheetnames
            sheet_pats = db_cfg.get("sheet_patterns", [])
            matched = [s for s in sheets if any(pat in s for pat in sheet_pats)]
            required = db_cfg.get("required_sheets", [])
            for rs in required:
                if rs not in matched:
                    errors.append(f"探采对比表中缺少「{rs}」工作表")
            db_cm = ColumnMatcher(db_cfg.get("patterns", {}))
            required_by_sheet = db_cfg.get("required_by_sheet", {})
            for sn in matched:
                required_cols = []
                for key_pat, cols in required_by_sheet.items():
                    if key_pat in sn:
                        required_cols = cols
                        break
                ws = wb[sn]
                rows = list(ws.iter_rows(min_row=1, values_only=True))
                if not rows:
                    continue
                r0, r1 = rows[0] if rows else [], rows[1] if len(rows) > 1 else []
                two_level = (r0 and r0[0] and is_header_row(r0[0]) and
                             r1 and r1[0] is None and any(c is not None for c in r1))
                if two_level:
                    parent = [str(c).strip() if c else None for c in r0]
                    child = [str(c).strip() if c else "" for c in r1]
                    last = ""
                    for i in range(len(parent)):
                        parent[i] = parent[i] if parent[i] else last
                        if parent[i]:
                            last = parent[i]
                    headers = [f"{p}_{c}" if p and c else (p or c) for p, c in zip(parent, child)]
                else:
                    h_idx = None
                    for i, r in enumerate(rows):
                        if r and r[0] and is_header_row(r[0]):
                            h_idx = i
                            break
                    if h_idx is None:
                        continue
                    headers = [str(c).strip() if c else "" for c in rows[h_idx]]
                for rc in required_cols:
                    if db_cm.match_header(headers, rc) is None:
                        suggestions = db_cm.suggest_headers(headers, rc)
                        sug_str = f"（建议检查是否包含：{', '.join(suggestions)}）" if suggestions else ""
                        errors.append(f"探采对比表「{sn}」缺少「{rc}」列{sug_str}")
            wb.close()
        except Exception as e:
            errors.append(f"读取探采对比表失败：{e}")

    return len(errors) == 0, errors

# ═══════════════════════════════════════════════════════════
# PPT 生成（不变）
# ═══════════════════════════════════════════════════════════

def generate_pptx(cfg, cfg_paths, png_files_map):
    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)

    DEEP_BLUE = PptRGBColor(0x2C, 0x3E, 0x50)
    ORANGE = PptRGBColor(0xE6, 0x7E, 0x22)
    WHITE = PptRGBColor(0xFF, 0xFF, 0xFF)
    LIGHT_GRAY = PptRGBColor(0xF2, 0xF3, 0xF5)
    BODY_TEXT = PptRGBColor(0x55, 0x55, 0x55)

    def _blank_slide():
        layout = prs.slide_layouts[6]
        return prs.slides.add_slide(layout)

    def _set_bg(slide, color):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _add_textbox(slide, left, top, width, height, text, font_size, color,
                     bold=False, alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = PptPt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = alignment
        return tf

    def _add_deco_strip(slide, left, top, width, height, color):
        strip = slide.shapes.add_shape(1, left, top, width, height)
        strip.fill.solid()
        strip.fill.fore_color.rgb = color
        strip.line.fill.background()
        return strip

    def _add_chart_image(slide, img_path, left, top, width, height):
        if os.path.exists(img_path):
            slide.shapes.add_picture(img_path, left, top, width, height)

    cover = _blank_slide()
    _set_bg(cover, DEEP_BLUE)
    _add_deco_strip(cover, PptInches(1.6), PptInches(3.2), PptInches(9.4), PptInches(0.05), ORANGE)
    _add_textbox(cover, PptInches(1.6), PptInches(1.4), PptInches(9.4), PptInches(1.8),
                 cfg.get('mine', {}).get('name', ''), 40, WHITE, bold=True)
    _add_textbox(cover, PptInches(1.6), PptInches(3.5), PptInches(9.4), PptInches(1.2),
                 f"{cfg_paths['year']}年{cfg_paths['cn_q_full']}探采对比分析报告", 26,
                 PptRGBColor(0xEC, 0xF0, 0xF1))
    _add_textbox(cover, PptInches(1.6), PptInches(5.0), PptInches(9.4), PptInches(0.8),
                 f"{cfg.get('mine', {}).get('location', '')}  |  {cfg_paths['report_date']}", 14,
                 PptRGBColor(0x95, 0xA5, 0xA6))

    chapters = [
        {"type": "section", "num": "1", "title": "对比依据与范围",
         "desc": "以ZHTK三维地质勘探模型为对比基准，以同区域生产勘探模型为验证对象"},
        {"type": "content", "title": "1  对比依据与范围",
         "text": "对比内容涵盖矿体形态（面积误差、重合率、歪曲率）、资源储量（矿石量、品位、金属量）以及开采过程中的矿石贫化与损失三项核心指标。",
         "imgs": None},
        {"type": "section", "num": "2", "title": "对比参数体系",
         "desc": "覆盖矿体几何形态与资源储量两个维度"},
        {"type": "section", "num": "3", "title": "探采对比结果分析",
         "desc": "矿体形态 + 资源储量 + 贫化损失 三项对比"},
        {"type": "content", "title": "3.1  矿体形态对比分析",
         "text": "面积重合率反映两模型矿体投影的重合程度，形态歪曲率表征矿体边界形态偏差。",
         "imgs": ["area_dr", "area_wr"]},
        {"type": "content", "title": "3.2  资源储量对比分析",
         "text": "对比两模型在各台阶的矿石量、品位、金属量，评估资源估算精度。",
         "imgs": ["ore_qty", "grade_line", "metal"]},
        {"type": "content", "title": "3.2  各台阶双模型对比",
         "text": "从矿石量、品位、金属量三个维度展示生产模型与ZHTK地质模型在各台阶的对比情况。",
         "imgs": ["step_ore", "step_grade", "step_metal"]},
        {"type": "section", "num": "4", "title": "结论与建议",
         "desc": "综合对比结论与生产建议"},
    ]

    for ch in chapters:
        if ch["type"] == "section":
            sec = _blank_slide()
            _set_bg(sec, DEEP_BLUE)
            _add_deco_strip(sec, PptInches(1.6), PptInches(3.4), PptInches(1.5), PptInches(0.05), ORANGE)
            _add_textbox(sec, PptInches(1.6), PptInches(2.2), PptInches(9.4), PptInches(1.0),
                         ch["num"], 48, ORANGE, bold=True)
            _add_textbox(sec, PptInches(1.6), PptInches(3.6), PptInches(9.4), PptInches(1.2),
                         ch["title"], 32, WHITE, bold=True)
            _add_textbox(sec, PptInches(1.6), PptInches(4.7), PptInches(9.4), PptInches(0.8),
                         ch["desc"], 16, PptRGBColor(0xBD, 0xC3, 0xC7))
        else:
            slide = _blank_slide()
            _set_bg(slide, LIGHT_GRAY)
            _add_deco_strip(slide, PptInches(0), PptInches(0), PptInches(0.08), PptInches(7.14), DEEP_BLUE)
            _add_deco_strip(slide, PptInches(0.35), PptInches(0.25), PptInches(12.0), PptInches(0.03), ORANGE)
            _add_textbox(slide, PptInches(0.5), PptInches(0.4), PptInches(11.5), PptInches(0.6),
                         ch["title"], 26, DEEP_BLUE, bold=True)
            _add_textbox(slide, PptInches(0.5), PptInches(1.0), PptInches(11.5), PptInches(0.5),
                         ch["text"], 13, BODY_TEXT)

            imgs = ch.get("imgs")
            if imgs:
                n = len(imgs)
                y_start = PptInches(1.65)
                if n == 2:
                    w_img, h_img = PptInches(5.8), PptInches(2.55)
                    positions = [(PptInches(0.55), y_start, w_img, h_img),
                                 (PptInches(6.5), y_start, w_img, h_img)]
                elif n == 3:
                    w_img, h_img = PptInches(3.85), PptInches(2.35)
                    positions = [(PptInches(0.55), y_start, w_img, h_img),
                                 (PptInches(4.55), y_start, w_img, h_img),
                                 (PptInches(8.55), y_start, w_img, h_img)]
                else:
                    positions = []
                for idx, key in enumerate(imgs):
                    if idx < len(positions) and key in png_files_map:
                        l, t, w, h = positions[idx]
                        _add_chart_image(slide, png_files_map[key], l, t, w, h)

    pptx_path = os.path.join(cfg_paths["output_dir"],
                             f"{cfg['mine']['name']}{cfg_paths['year']}年{cfg_paths['cn_q_full']}探采对比分析_演示版.pptx")
    prs.save(pptx_path)
    return pptx_path

# ═══════════════════════════════════════════════════════════
# 核心报告生成（使用真正解耦的图表工厂）
# ═══════════════════════════════════════════════════════════

def make_report(cfg_paths, cfg, ph_data, db_data):
    rcfg = cfg.get("report", {})
    grading = cfg.get("grading", {})
    analyze = cfg.get("analyze", {})
    col_match = ColumnMatcher(cfg["columns"]["duibai"]["patterns"])
    cfg_styles = cfg["styles"]

    y = cfg_paths["year"]
    q = cfg_paths["quarter"]
    cn_q_full = cfg_paths["cn_q_full"]
    CN_Q = cfg_paths["CN_Q"]
    report_date = cfg_paths["report_date"]
    month_labels = cfg_paths["month_labels"]
    month_sheets = cfg_paths["month_sheets"]
    quarter_sheet_name = cfg_paths["quarter_sheet_name"]

    output_dir = cfg_paths["output_dir"]
    os.makedirs(output_dir, exist_ok=True)
    png_files_map = {}

    # ===== V10 特性：SQLite 历史数据库 =====
    db_path = os.path.join(output_dir, "probe_history.db")
    init_database(db_path)

    # V12 优化：render_and_save 改为先存 PNG 再用文件路径插入 Word，只渲染一次
    def render_and_save(chart_func, args, img_width, fig_caption, png_key):
        """
        V12 优化：调用图表工厂函数获取 (fig, buf)，先将 fig 保存为 PNG，
        再用 PNG 文件路径插入 Word（doc.add_picture 支持文件路径），
        同一张图只渲染一次，buf 不再用于 Word 插入。
        """
        fig, buf = chart_func(*args)          # 图表工厂负责渲染
        # V12 优化：先保存 PNG 到文件
        png_path = os.path.join(output_dir, f"{png_key}.png")
        fig.savefig(png_path, dpi=200, bbox_inches='tight', facecolor='white')
        # V12 优化：用文件路径插入 Word（不再使用 buf）
        doc.add_picture(png_path, width=Inches(img_width))
        add_para(doc, fig_caption, 'MiningNoIndent')
        plt.close(fig)
        png_files_map[png_key] = png_path
        return png_path

    doc = Document()
    setup_document_styles(doc, cfg_styles)

    # 封面
    for _ in range(7):
        doc.add_paragraph()
    add_para(doc, f"{cfg['mine']['name']}\n{y}年{cn_q_full}探采对比分析报告", 'MiningTitle')
    doc.add_paragraph()
    add_para(doc, cfg["mine"].get("subtitle", ""), 'MiningSubtitle')
    doc.add_paragraph()
    doc.add_paragraph()
    add_para(doc, f"{cfg['mine'].get('location', '')}\n{report_date}", 'MiningSubtitle')
    doc.add_page_break()

    # 第1章
    add_heading_styled(doc, "1  对比依据与范围", 1)
    add_para(doc, rcfg.get("chapter1_text", ""))
    add_para(doc, "对比内容涵盖矿体形态（面积误差、重合率、歪曲率）、资源储量（矿石量、品位、金属量）以及开采过程中的矿石贫化与损失三项核心指标。")

    # 第2章
    add_heading_styled(doc, "2  对比参数体系", 1)
    add_para(doc, "探采对比选取以下核心参数，覆盖矿体几何形态与资源储量两个维度：")
    params = [
        ("面积绝对误差 Sδ", "Sδ = Su - Sc", "ZHTK面积与生产面积之差（m²）"),
        ("面积相对误差 Sr", "Sr = (Su - Sc) / Su × 100%", "面积偏差相对于ZHTK模型的比例"),
        ("面积重合率 Dr", "Dr = Sd / Su × 100%", "两模型重合部分占ZHTK模型面积比"),
        ("形态歪曲率 Wr", "Wr = (Sn+Sp) / Su × 100%", "不重合部分占总面积的比例"),
        ("矿石量相对误差 Qr", "Qr = (Qu - Qc) / Qu × 100%", "矿石量偏差相对于ZHTK模型的比例"),
        ("品位绝对误差 Cδ", "Cδ = Cu - Cc", "ZHTK品位与生产品位的差值（百分点）"),
        ("金属量相对误差 Pr", "Pr = (Pu - Pc) / Pu × 100%", "金属量偏差相对于ZHTK模型的比例"),
    ]
    make_table(doc, ["参数名称", "计算公式", "说明"], [[n, f, d] for n, f, d in params], cfg_styles)
    add_para(doc, "表2-1  探采对比参数一览", 'MiningNoIndent')

    # 第3章
    add_heading_styled(doc, "3  探采对比结果分析", 1)

    # 3.1 矿体形态
    add_heading_styled(doc, "3.1  矿体形态对比分析", 2)
    area_recs = db_data.get("面积误差重合率", [])
    step_area = [r for r in area_recs if is_step_row(r)]
    area_rows = []
    sum_prod = sum_zhtk = sum_overlap = sum_wadd = sum_wsub = 0.0
    step_wr_map = {}
    step_dr_map = {}
    step_area_data = {}

    for rec in step_area:
        step = rec.get("_step", "")
        prod = sf(col_match.get_or(rec, "面积_生产"))
        zhtk = sf(col_match.get_or(rec, "面积_ZHTK"))
        if not zhtk:
            continue
        abs_e = zhtk - prod
        rel_e = f"{(abs_e/zhtk*100):+.0f}%" if zhtk else "-"
        over = sf(col_match.get_or(rec, "面积_重合"))
        over_r = f"{over/zhtk*100:.0f}%" if zhtk else "-"
        wadd = sf(col_match.get_or(rec, "面积_歪曲增加"))
        wsub = sf(col_match.get_or(rec, "面积_歪曲减少"))
        wttl = wadd + wsub
        wr_str = f"{wttl/zhtk*100:.0f}%" if zhtk else "-"
        area_rows.append([step, fm(prod, 2), fm(zhtk, 2), fm(abs_e, 2), rel_e,
                          fm(over, 2), over_r, fm(wadd, 2), fm(wsub, 2), fm(wttl, 2), wr_str])
        sum_prod += prod
        sum_zhtk += zhtk
        sum_overlap += over
        sum_wadd += wadd
        sum_wsub += wsub
        if zhtk:
            step_wr_map[step] = wttl / zhtk * 100
            step_dr_map[step] = over / zhtk * 100
            step_area_data[step] = {
                "prod": prod, "zhtk": zhtk, "over": over,
                "wadd": wadd, "wsub": wsub,
                "dr": over / zhtk * 100, "wr": wttl / zhtk * 100,
                "rel_e": abs_e / zhtk * 100
            }

    if sum_zhtk:
        t_abs = sum_zhtk - sum_prod
        t_rel = f"{(t_abs/sum_zhtk*100):+.0f}%"
        t_over_r = f"{sum_overlap/sum_zhtk*100:.0f}%"
        t_w = sum_wadd + sum_wsub
        t_wr = f"{t_w/sum_zhtk*100:.0f}%"
        area_rows.append(["合计", fm(sum_prod, 2), fm(sum_zhtk, 2), fm(t_abs, 2), t_rel,
                          fm(sum_overlap, 2), t_over_r, fm(sum_wadd, 2), fm(sum_wsub, 2), fm(t_w, 2), t_wr])
    else:
        t_rel = t_over_r = t_wr = "-"

    add_para(doc, f"ZHTK模型矿体总面积{fm(sum_zhtk,0)}m²，生产模型矿体总面积{fm(sum_prod,0)}m²。面积重合率{t_over_r}，形态歪曲率{t_wr}。")

    if sum_zhtk:
        ovr_pct = sum_overlap / sum_zhtk * 100
        if ovr_pct >= 80:
            morph_desc = '较高（≥80%），多数台阶矿体形态吻合良好'
        elif ovr_pct >= 60:
            morph_desc = '中等（60%~80%），部分台阶存在局部形态偏差，但整体可接受'
        else:
            morph_desc = '偏低（<60%），两个模型在矿体空间形态上存在较大差异，需逐台阶复核矿体边界圈定'
        add_para(doc, f"从面积重合率来看，两个模型在矿体空间展布上的整体一致性{morph_desc}。")
        t_wr_num = t_w / sum_zhtk * 100 if sum_zhtk else 0
        if t_wr_num > 30:
            wr_desc = '超出30%警戒线，矿体形态在两个模型中偏差明显，可能影响资源储量估算精度，建议技术组重点复核淋滤层剔除及岩粉样化验系统误差'
        elif t_wr_num > 20:
            wr_desc = '处于20%~30%区间，存在一定形态偏差，部分台阶矿体圈定需关注'
        else:
            wr_desc = '处于可接受范围内（≤20%），矿体形态在两个模型中一致性良好，矿体边界圈定精度较高'
        add_para(doc, f"形态歪曲率{t_wr}，{wr_desc}。")
        area_chart_steps = [r.get("_step", "") for r in step_area if r.get("_step", "") != "合计"]
        if area_chart_steps:
            dr_vals = [step_area_data.get(s, {}).get("dr", 0) for s in area_chart_steps]
            wr_vals = [step_area_data.get(s, {}).get("wr", 0) for s in area_chart_steps]
            high_dr = [s for i, s in enumerate(area_chart_steps) if dr_vals[i] >= 80]
            low_dr = [s for i, s in enumerate(area_chart_steps) if dr_vals[i] < 60]
            high_wr = [s for i, s in enumerate(area_chart_steps) if wr_vals[i] > 30]
            parts = []
            if high_dr:
                parts.append(f"面积重合率达到A级（≥80%）的台阶包括{'、'.join(high_dr)}")
            if low_dr:
                parts.append(f"面积重合率偏低（<60%）的台阶包括{'、'.join(low_dr)}，建议重点复核")
            if high_wr:
                parts.append(f"形态歪曲率超警戒线（>30%）的台阶包括{'、'.join(high_wr)}，矿体空间形态差异明显")
            if not parts:
                parts.append("各台阶面积重合率均处于合理区间，未发现显著异常台阶")
            add_para(doc, "；".join(parts) + "。")
        else:
            add_para(doc, "各台阶面积数据待补充，暂不做逐台阶分析。")
    else:
        add_para(doc, "矿体面积数据暂缺，无法进行形态对比分析。")

    doc.add_paragraph()
    render_and_save(make_area_dr_chart, (area_chart_steps, step_area_data, cfg),
                    6.5, "图3-1  各台阶矿体面积及重合面积组合图", "area_dr")
    doc.add_paragraph()
    render_and_save(make_area_wr_chart, (area_chart_steps, step_area_data, cfg),
                    6.5, "图3-2  各台阶歪曲面积及歪曲率组合图", "area_wr")

    # 3.2 资源储量
    add_heading_styled(doc, "3.2  资源储量对比分析", 2)
    ore_recs = db_data.get("矿量对比", [])
    ore_steps = [r for r in ore_recs if is_step_row(r)]
    ore_sums = [r for r in ore_recs if not is_step_row(r)]
    ore_all = ore_steps + ore_sums

    n_fixed = validate_ore_logic(ore_all, col_match, "矿量对比-")
    if n_fixed:
        print(f"  共修正 {n_fixed} 处金属量计算尾差")

    step_detail = {}
    for rec in ore_steps:
        nm = rec.get("_step", "")
        p_q = sf(col_match.get_or(rec, "生产_矿石量"))
        p_c = sf(col_match.get_or(rec, "生产_品位"))
        p_p = sf(col_match.get_or(rec, "生产_金属量"))
        z_q = sf(col_match.get_or(rec, "ZHTK_矿石量"))
        z_c = sf(col_match.get_or(rec, "ZHTK_品位"))
        z_p = sf(col_match.get_or(rec, "ZHTK_金属量"))
        step_detail[nm] = {
            "p_q": p_q, "p_c": p_c, "p_p": p_p,
            "z_q": z_q, "z_c": z_c, "z_p": z_p,
            "q_diff": p_q - z_q,
            "q_diff_pct": (p_q - z_q) / z_q * 100 if z_q else 0
        }

    total = next((r for r in ore_all if r.get("_step", "") == "合计"), None)
    if total:
        t_q = sf(col_match.get_or(total, "生产_矿石量"))
        t_c = sf(col_match.get_or(total, "生产_品位"))
        t_p = sf(col_match.get_or(total, "生产_金属量"))
        z_q = sf(col_match.get_or(total, "ZHTK_矿石量"))
        z_c = sf(col_match.get_or(total, "ZHTK_品位"))
        z_p = sf(col_match.get_or(total, "ZHTK_金属量"))

    # ===== V10 特性：保存历史数据并异常检测 =====
    save_indicators_to_db(db_path, y, q, step_detail, {
        "t_q": t_q, "z_q": z_q, "t_c": t_c, "z_c": z_c, "t_p": t_p, "z_p": z_p
    })
    is_anomaly, change_pct, prev_val = detect_anomaly(db_path, y, q, "ore", threshold=cfg.get("analyze", {}).get("anomaly_threshold", 30))
    if is_anomaly:
        warnings.warn(f"矿石量异常：环比变动 {change_pct:+.1f}%，上季度值 {prev_val}")

    add_heading_styled(doc, "3.2.1  双模型矿石量与金属量统计", 3)
    if total:
        add_para(doc, f"生产勘探模型：全矿区矿石量 {fm(t_q,0)} t，平均品位 {t_c:.3f}%，金属量 {fm(t_p,0)} t。")
        add_para(doc, f"ZHTK地质勘探模型：全矿区矿石量 {fm(z_q,0)} t，平均品位 {z_c:.3f}%，金属量 {fm(z_p,0)} t。")
        add_para(doc, f"矿石量绝对偏差 {fm(abs(t_q-z_q),0)} t（偏差率 {(t_q-z_q)/z_q*100:+.1f}%），金属量总偏差 {fm(abs(t_p-z_p),0)} t（偏差率 {(t_p-z_p)/z_p*100:+.1f}%）。")

    add_heading_styled(doc, "3.2.2  矿石量变化分析", 3)

    pos_steps = [s for s, d in step_detail.items() if d["q_diff"] > 0]
    neg_steps = [s for s, d in step_detail.items() if d["q_diff"] < 0]
    pos_steps.sort(key=lambda s: abs(step_detail[s]["q_diff"]), reverse=True)
    neg_steps.sort(key=lambda s: abs(step_detail[s]["q_diff"]), reverse=True)

    if pos_steps:
        add_heading_styled(doc, "（一）矿石量正变台阶", 3)
        add_para(doc, f"矿石量正变（生产模型矿石量 > ZHTK模型矿石量）台阶共 {len(pos_steps)} 个，即生产勘探模型在这些台阶揭示的矿体规模大于ZHTK地质勘探模型的预估值。正变台阶合计矿石量正变 {fm(sum(step_detail[s]['q_diff'] for s in pos_steps),0)} t。")
        for step in pos_steps:
            d = step_detail[step]
            parts = []
            parts.append(f"{step}台阶：生产矿石量 {fm(d['p_q'],0)} t（品位 {d['p_c']:.3f}%），ZHTK矿石量 {fm(d['z_q'],0)} t（品位 {d['z_c']:.3f}%），矿石量正变 {fm(abs(d['q_diff']),0)} t（正变率 {d['q_diff_pct']:+.1f}%）")
            if d["p_c"] and d["z_c"]:
                gc = d["p_c"] - d["z_c"]
                parts.append(f"品位偏差 {gc:+.3f} 个百分点")
            add_para(doc, "；".join(parts) + "。")
    else:
        add_para(doc, "（一）矿石量正变台阶：无。")

    if neg_steps:
        add_heading_styled(doc, "（二）矿石量负变台阶", 3)
        add_para(doc, f"矿石量负变（生产模型矿石量 < ZHTK模型矿石量）台阶共 {len(neg_steps)} 个，即生产勘探模型在这些台阶揭示的矿体规模小于ZHTK地质勘探模型的预估值。负变台阶合计矿石量负变 {fm(sum(abs(step_detail[s]['q_diff']) for s in neg_steps),0)} t。")
        for step in neg_steps:
            d = step_detail[step]
            parts = []
            parts.append(f"{step}台阶：生产矿石量 {fm(d['p_q'],0)} t（品位 {d['p_c']:.3f}%），ZHTK矿石量 {fm(d['z_q'],0)} t（品位 {d['z_c']:.3f}%），矿石量负变 {fm(abs(d['q_diff']),0)} t（负变率 {abs(d['q_diff_pct']):.1f}%）")
            if d["p_c"] and d["z_c"]:
                gc = d["p_c"] - d["z_c"]
                parts.append(f"品位偏差 {gc:+.3f} 个百分点")
            add_para(doc, "；".join(parts) + "。")
    else:
        add_para(doc, "（二）矿石量负变台阶：无。")

    ore_chart_steps = auto_sort_steps([s for s, d in step_detail.items() if d["q_diff_pct"] != 0 or d["p_q"] > 0 or d["z_q"] > 0])

    doc.add_paragraph()
    render_and_save(make_ore_qty_chart, (ore_chart_steps, step_detail, cfg),
                    6.5, "图3-3  各台阶矿石量对比柱状图", "ore_qty")
    doc.add_paragraph()
    render_and_save(make_grade_line_chart, (ore_chart_steps, step_detail, cfg),
                    6.2, "图3-4  各台阶品位对比折线图", "grade_line")
    doc.add_paragraph()
    render_and_save(make_metal_chart, (ore_chart_steps, step_detail, cfg),
                    6.5, "图3-5  各台阶金属量对比柱状图", "metal")

    # 跨 Sheet 台阶一致性校验
    db_ore_steps = set(step_detail.keys())
    db_area_steps = set()
    for r in db_data.get("矿量对比", []):
        if is_step_row(r):
            db_area_steps.add(r.get("_step", ""))
    if db_ore_steps != db_area_steps:
        warnings.warn(f"台阶列表不一致: 仅矿量对比表有 {db_ore_steps - db_area_steps}, "
                      f"仅面积表有 {db_area_steps - db_ore_steps}")

    doc.add_paragraph()
    render_and_save(make_step_ore_chart, (ore_chart_steps, step_detail, cfg),
                    6.5, "图3-6  各台阶矿石量对比柱状图", "step_ore")
    doc.add_paragraph()
    render_and_save(make_step_grade_chart, (ore_chart_steps, step_detail, cfg),
                    6.2, "图3-7  各台阶品位对比折线图", "step_grade")
    doc.add_paragraph()
    render_and_save(make_step_metal_chart, (ore_chart_steps, step_detail, cfg),
                    6.5, "图3-8  各台阶金属量对比柱状图", "step_metal")

    # 3.3 贫化损失
    add_heading_styled(doc, "3.3  矿石贫化与损失分析", 2)
    pinghua_col = ColumnMatcher(cfg["columns"]["pinghua"]["patterns"])
    q_data = ph_data.get(quarter_sheet_name, {})
    m_data = {ms: ph_data.get(ms, {}) for ms in month_sheets}

    all_steps = set()
    for d in [*m_data.values(), q_data]:
        all_steps.update(k for k in d if k.strip() and k not in ("备注", "签名", "说明"))
    order = rcfg.get("step_order", [])
    steps = [s for s in order if s in all_steps] + [s for s in sorted(all_steps) if s not in order]

    def get_ph_val(data_dict, step, field):
        rec = data_dict.get(step, {})
        return pinghua_col.get_or(rec, field, "-")

    tbl_header = ["台阶"]
    for ml in month_labels:
        tbl_header.append(f"{ml}贫化率")
        tbl_header.append(f"{ml}损失率")
    tbl_header.append(f"{CN_Q}季度贫化率")
    tbl_header.append(f"{CN_Q}季度损失率")

    ph_rows = []
    for step in steps:
        row = [step]
        for md in list(m_data.values()) + [q_data]:
            ph_val = get_ph_val(md, step, "贫化率")
            lo_val = get_ph_val(md, step, "损失率")
            row.append(fp(ph_val) if isinstance(ph_val, (int, float)) and ph_val != "-" else str(ph_val) if ph_val else "-")
            row.append(fp(lo_val) if isinstance(lo_val, (int, float)) and lo_val != "-" else str(lo_val) if lo_val else "-")
        ph_rows.append(row)

    make_table(doc, tbl_header, ph_rows, cfg_styles)
    add_para(doc, f"表3-3  {y}年{cn_q_full}矿石贫化损失统计表", 'MiningNoIndent')

    # 第4章
    add_heading_styled(doc, "4  结论与建议", 1)

    actual_mining = 0.0
    actual_grade = 0.0
    if q_data:
        tqr = q_data.get("合计", {})
        actual_mining = sf(pinghua_col.get_or(tqr, "出矿量", 0))
        actual_grade = sf(pinghua_col.get_or(tqr, "品位", 0))
    total_q = q_data.get("合计", {})
    t_ph = pinghua_col.get_or(total_q, "贫化率", "-")
    t_lo = pinghua_col.get_or(total_q, "损失率", "-")

    conclusions = []

    if total:
        conclusions.append(
            f"（1）双模型矿石量与金属量统计：生产勘探模型全矿区矿石量 {fm(t_q,0)} t（平均品位 {t_c:.3f}%），金属量 {fm(t_p,0)} t；"
            f"ZHTK地质勘探模型全矿区矿石量 {fm(z_q,0)} t（平均品位 {z_c:.3f}%），金属量 {fm(z_p,0)} t。"
            f"矿石量偏差率 {(t_q-z_q)/z_q*100:+.1f}%，金属量偏差率 {(t_p-z_p)/z_p*100:+.1f}%。"
        )
    else:
        conclusions.append(f"（1）双模型矿石量与金属量统计：数据暂缺。")

    if pos_steps and neg_steps:
        pos_total = sum(step_detail[s]["q_diff"] for s in pos_steps)
        neg_total = sum(abs(step_detail[s]["q_diff"]) for s in neg_steps)
        res_str = (
            f"（2）资源量对比分析：本次探采对比共涉及 {len(ore_steps)} 个台阶，"
            f"其中矿石量正变台阶 {len(pos_steps)} 个，合计正变 {fm(pos_total,0)} t；"
            f"负变台阶 {len(neg_steps)} 个，合计负变 {fm(neg_total,0)} t。"
            f"整体矿石量偏差率 {(t_q-z_q)/z_q*100:+.1f}%。"
        )
        all_devs = sorted(step_detail.items(), key=lambda x: abs(x[1]["q_diff"]), reverse=True)
        top_steps = all_devs[:3]
        top_desc = "、".join([f"{s[0]}台阶段（{'正变' if s[1]['q_diff']>0 else '负变'}{fm(abs(s[1]['q_diff']),0)} t）" for s in top_steps])
        res_str += f"对矿石量总偏差影响最大的台阶依次为：{top_desc}。"
        conclusions.append(res_str)
    elif pos_steps:
        pos_total = sum(step_detail[s]["q_diff"] for s in pos_steps)
        conclusions.append(
            f"（2）资源量对比分析：{len(ore_steps)} 个台阶全部呈现矿石量正变，合计正变 {fm(pos_total,0)} t，"
            f"生产勘探模型在各台阶揭示的矿体规模均大于ZHTK地质勘探模型预估值。"
        )
    elif neg_steps:
        neg_total = sum(abs(step_detail[s]["q_diff"]) for s in neg_steps)
        conclusions.append(
            f"（2）资源量对比分析：{len(ore_steps)} 个台阶全部呈现矿石量负变，合计负变 {fm(neg_total,0)} t，"
            f"ZHTK地质勘探模型在各台阶均高估了矿体规模。"
        )

    main_steps_for_conc = sorted(
        [(s, d) for s, d in step_detail.items() if abs(d["q_diff"]) > 0],
        key=lambda x: abs(x[1]["q_diff"]), reverse=True
    )[:5]
    if main_steps_for_conc:
        step_desc_parts = []
        for s, d in main_steps_for_conc:
            direction = "正变" if d["q_diff"] > 0 else "负变"
            desc = f"{s}台阶{direction}{fm(abs(d['q_diff']),0)} t（偏差率{abs(d['q_diff_pct']):.1f}%）"
            step_desc_parts.append(desc)
        conc4 = f"（3）主要生产台阶对比分析：矿石量偏差最大的前 {len(main_steps_for_conc)} 个台阶为——{'；'.join(step_desc_parts)}。"
        max_dev = main_steps_for_conc[0]
        if max_dev[1]["q_diff"] > 0:
            conc4 += f"其中{max_dev[0]}台阶正变量最大，建议核实该台阶是否存在矿体边界外扩或品位估值偏高的情况。"
        else:
            conc4 += f"其中{max_dev[0]}台阶负变量最大，建议核实该台阶是否存在矿体尖灭或品位估值偏低的情况。"
        conclusions.append(conc4)
    else:
        conclusions.append(f"（3）主要生产台阶对比分析：各台阶矿石量偏差均在合理范围内。")

    if actual_mining > 0:
        mining_diff = actual_mining - t_q
        conc5 = (
            f"（4）实际生产情况与生产模型对比：{cn_q_full}实际采矿量 {fm(actual_mining,0)} t（品位 {actual_grade:.3f}%），"
            f"较生产勘探模型{'减少' if mining_diff < 0 else '增加'}{fm(abs(mining_diff),0)} t。"
            f"综合贫化率 {fp(t_ph) if isinstance(t_ph,(int,float)) else t_ph}，综合损失率 {fp(t_lo) if isinstance(t_lo,(int,float)) else t_lo}。"
        )
        if mining_diff < 0 and t_q and abs(mining_diff) / t_q > 0.05:
            conc5 += "实际采矿量低于生产模型估值超过5%，需关注开采过程中的矿量损失。"
        elif mining_diff > 0 and t_q and mining_diff / t_q > 0.05:
            conc5 += "实际采矿量高于生产模型估值超过5%，需关注是否存在超采或贫化控制不足。"
        conclusions.append(conc5)
    else:
        conclusions.append(
            f"（4）生产指标：综合贫化率 {fp(t_ph) if isinstance(t_ph,(int,float)) else t_ph}，综合损失率 {fp(t_lo) if isinstance(t_lo,(int,float)) else t_lo}。"
        )

    conclusions.append(
        f"（5）数据逻辑校验：已对 Excel 原始数据执行矿石量×品位=金属量的勾稽校验，发现并自动修正 {n_fixed} 处计算尾差，确保报告数据数学完备性。"
    )

    for c in conclusions:
        add_para(doc, c)
        doc.add_paragraph()

    docx_path = os.path.join(output_dir,
                             f"{cfg['mine']['name']}{y}年{cn_q_full}探采对比分析_生成版.docx")
    doc.save(docx_path)

    pptx_path = generate_pptx(cfg, cfg_paths, png_files_map)
    return docx_path, pptx_path, png_files_map

# ═══════════════════════════════════════════════════════════
# 文件查找与命令行
# ═══════════════════════════════════════════════════════════

def find_file_by_keywords(base_dir, keywords, ext=".xlsx"):
    target_dir = os.path.join(base_dir, "附表")
    if not os.path.isdir(target_dir):
        return None
    candidates = []
    kw_lower = [k.lower() for k in keywords]
    for fname in os.listdir(target_dir):
        if not fname.lower().endswith(ext.lower()):
            continue
        stem = fname.lower()
        if all(k in stem for k in kw_lower):
            candidates.append(os.path.join(target_dir, fname))
    if not candidates:
        return None
    candidates.sort(key=lambda p: sum(1 for kw in kw_lower if kw in os.path.basename(p).lower()), reverse=True)
    return candidates[0]

def construct_file_path(base_dir, pattern, year, month_range, cn_q):
    fname = pattern.format(year=year, month_range=month_range, cn_q=cn_q)
    path = os.path.join(base_dir, "附表", fname)
    if os.path.exists(path):
        return path
    return None

def parse_args():
    parser = argparse.ArgumentParser(description="矿山探采对比分析报告生成工具 v12.0")
    parser.add_argument("--year", type=int, help="年份（如 2026）")
    parser.add_argument("--quarter", type=int, choices=[1,2,3,4], help="季度（1/2/3/4）")
    parser.add_argument("--pinghua", help="贫化损失统计表路径（直接指定，跳过自动查找）")
    parser.add_argument("--duibai", help="探采对比表格路径（直接指定，跳过自动查找）")
    parser.add_argument("--config", default="config.json", help="配置文件路径（默认 config.json）")
    parser.add_argument("--output", help="输出目录（默认桌面）")
    return parser.parse_args()

# ═══════════════════════════════════════════════════════════
# 主程序
# ═══════════════════════════════════════════════════════════

def main():
    args = parse_args()

    # 加载配置
    cfg = load_config(args.config)

    # 确定年份和季度（交互式 + 命令行双重模式）
    if args.year and args.quarter:
        # 命令行模式：静默使用传入参数
        year = args.year
        quarter = args.quarter
    else:
        # 交互模式：弹窗让用户输入
        print("\n【交互模式】请确认要分析的数据时间：")
        now = datetime.now()
        default_year = args.year if args.year else now.year
        default_quarter = args.quarter if args.quarter else (now.month - 1) // 3 + 1

        year_input = input(f"请输入年份（直接回车默认 {default_year}）：").strip()
        year = int(year_input) if year_input else default_year

        quarter_input = input(f"请输入季度（1/2/3/4，直接回车默认 {default_quarter}）：").strip()
        quarter = int(quarter_input) if quarter_input else default_quarter

        # 合法性校验
        if quarter not in [1, 2, 3, 4]:
            print(f"⚠ 季度输入有误（{quarter}），已自动调整为 {default_quarter}")
            quarter = default_quarter
        print(f"\n▶ 即将处理：{year} 年第 {quarter} 季度数据")

    if quarter not in [1, 2, 3, 4]:
        quarter = 1

    CN_Q = {1: "一", 2: "二", 3: "三", 4: "四"}[quarter]
    cn_q_full = f"第{CN_Q}季度"
    yy = str(year)[-2:]
    ms, me = (quarter - 1) * 3 + 1, quarter * 3
    month_range = f"{ms:02d}-{me:02d}"
    month_labels = [f"{m}月" for m in range(ms, me + 1)]
    month_sheets = [f"{yy}-{m:02d}" for m in range(ms, me + 1)]
    quarter_sheet_name = "一季度" if CN_Q == "一" else f"{CN_Q}季度"

    rp_m, rp_y = me + 1, year + (1 if me == 12 else 0)
    if rp_m > 12:
        rp_m = 1
    report_date = f"{rp_y}年{rp_m}月"

    # 输出目录
    if args.output:
        output_dir = args.output
    else:
        desktop = os.path.join(os.path.expanduser("~"), "Desktop")
        output_dir = os.path.join(desktop, f"{year}年第{CN_Q}季度探采对比分析输出")
    os.makedirs(output_dir, exist_ok=True)

    script_dir = os.path.dirname(os.path.abspath(__file__))

    # 查找/确认文件
    f_pinghua = args.pinghua if args.pinghua else None
    f_duibai = args.duibai if args.duibai else None

    if not f_pinghua:
        pattern = cfg.get("files", {}).get("pinghua_pattern", DEFAULT_CONFIG["files"]["pinghua_pattern"])
        f_pinghua = construct_file_path(script_dir, pattern, year, month_range, CN_Q)
        if not f_pinghua:
            keywords = cfg.get("files", {}).get("pinghua_keywords", DEFAULT_CONFIG["files"]["pinghua_keywords"])
            f_pinghua = find_file_by_keywords(script_dir, keywords)
            if f_pinghua:
                print(f"✓ 自动匹配贫化损失表：{os.path.basename(f_pinghua)}")

    if not f_duibai:
        pattern = cfg.get("files", {}).get("duibai_pattern", DEFAULT_CONFIG["files"]["duibai_pattern"])
        f_duibai = construct_file_path(script_dir, pattern, year, month_range, CN_Q)
        if not f_duibai:
            keywords = cfg.get("files", {}).get("duibai_keywords", DEFAULT_CONFIG["files"]["duibai_keywords"])
            f_duibai = find_file_by_keywords(script_dir, keywords)
            if f_duibai:
                print(f"✓ 自动匹配探采对比表：{os.path.basename(f_duibai)}")

    if not f_pinghua or not os.path.exists(f_pinghua):
        print(f"✗ 找不到贫化损失统计表（尝试路径：{f_pinghua}）")
        print("  请使用 --pinghua 指定文件路径，或确保文件在「附表」目录中。")
        sys.exit(1)

    if not f_duibai or not os.path.exists(f_duibai):
        print(f"✗ 找不到探采对比表格（尝试路径：{f_duibai}）")
        print("  请使用 --duibai 指定文件路径，或确保文件在「附表」目录中。")
        sys.exit(1)

    cfg_paths = {
        "year": year,
        "quarter": quarter,
        "yy": yy,
        "CN_Q": CN_Q,
        "cn_q_full": cn_q_full,
        "month_range": month_range,
        "month_sheets": month_sheets,
        "month_labels": month_labels,
        "report_date": report_date,
        "quarter_sheet_name": quarter_sheet_name,
        "base_dir": script_dir,
        "f_pinghua": f_pinghua,
        "f_duibai": f_duibai,
        "output_dir": output_dir,
    }

    print("\n" + "=" * 64)
    print(f"  {cfg['mine']['name']} 探采对比分析报告生成工具 v12.0（优化版）")
    print("  性能优化 + Bug 修复 + 代码质量提升 + V10 智能识别 + 同义词 80+")
    print("=" * 64)
    print(f"\n  年份：{year}  季度：{quarter}（{cn_q_full}）")
    print(f"  贫化损失表：{f_pinghua}")
    print(f"  探采对比表：{f_duibai}")
    print(f"  输出目录：{output_dir}")

    print(f"\n▶ 正在预检文件结构...")
    passed, errors = verify_excel_structure(f_pinghua, f_duibai, cfg)
    if not passed:
        print(f"\n  共发现 {len(errors)} 个问题：")
        for e in errors:
            print(f"    ✗ {e}")
        print("\n【提示】请根据上述提示调整 Excel 表头或修改 config.json 中的列名匹配规则。")
        sys.exit(1)
    print("  预检通过 ✓")

    print(f"\n▶ 正在读取数据...")
    ph_data = read_pinghua(f_pinghua, cfg)
    db_data = read_duibai(f_duibai, cfg)
    print(f"  已读取 {len(ph_data)} 张贫化损失表，{len(db_data)} 组探采对比表")

    print(f"\n▶ 正在生成报告（Word+PPT+图表PNG，请稍候）...")
    docx_path, pptx_path, png_files_map = make_report(cfg_paths, cfg, ph_data, db_data)

    print(f"\n{'='*60}")
    print(f"✅ 全部生成完成！")
    print(f"  Word 报告：{docx_path}")
    print(f"  PPT 演示：{pptx_path}")
    print(f"  图表 PNG：{len(png_files_map)} 张")
    print(f"\n  所有文件已保存到：{output_dir}")

if __name__ == "__main__":
    try:
        main()
    except KeyboardInterrupt:
        print("\n\n操作已取消。")
        sys.exit(0)
    except Exception as e:
        import traceback
        print("\n" + "=" * 64)
        print("【程序运行出错】")
        print(f"错误类型：{type(e).__name__}")
        print(f"错误信息：{e}")
        print("-" * 64)
        traceback.print_exc()
        print("-" * 64)
        print("常见问题排查：")
        print("  1. 请确认「附表」文件夹存在且 Excel 文件未损坏")
        print("  2. 请确认 config.json 中的列名关键字与 Excel 表头匹配")
        print("  3. 请确认已安装 python-docx、openpyxl、matplotlib、python-pptx")
        sys.exit(1)
